from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, layout_index, title_text, content_text=None):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    if title_text:
        title = slide.shapes.title
        title.text = title_text
    
    # Set Content
    if content_text and layout_index != 0: # Skip content for title slide usually
        # Find the content placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                text_frame = shape.text_frame
                text_frame.text = content_text
                break
    
    return slide

def add_text_slide_with_bullets(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear() # Clear existing paragraphs
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

prs = Presentation()

# --- Slide 1: Introduction ---
slide_1 = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Introduction to RAG and Its Application"
subtitle.text = "Presented By: Aayush Srivastava\nNash Tech."

# --- Slide 2: Agenda ---
agenda_points = [
    "1. Introduction: What is LLM, What is RAG?",
    "2. LLM And Its Limitations: Why RAG is important?",
    "3. RAG Architecture",
    "4. How Does RAG Work",
    "5. RAG Vs Fine-Tuning",
    "6. Benefits Of RAG",
    "7. Applications",
    "8. Demo"
]
add_text_slide_with_bullets(prs, "Agenda", agenda_points)

# --- Slide 3: Introduction (Section Header) ---
add_slide(prs, 2, "Introduction") # Title Only

# --- Slide 4: What is LLM ---
llm_points = [
    "A large language model (LLM) is a type of artificial intelligence program that can recognize and generate text, among other tasks.",
    "LLMs are very large models pre-trained on vast amounts of data.",
    "Built on transformer architecture: a neural network with an encoder and decoder having self-attention capabilities.",
    "Tasks: Answering questions, summarizing documents, translating languages, completing sentences.",
    "Example: OpenAI's GPT-3 has 175 billion parameters and supports inputs up to 100K tokens."
]
add_text_slide_with_bullets(prs, "What is LLM", llm_points)

# --- Slide 5: What is LLM (Continued) ---
llm_cont_points = [
    "In simpler terms, an LLM is a computer program fed enough examples to recognize and interpret human language or complex data.",
    "Quality of samples impacts learning; programmers often use curated datasets.",
    "[PLACEHOLDER FOR DIAGRAM: Generative Approach / Model Architecture]"
]
add_text_slide_with_bullets(prs, "What is LLM (Continued)", llm_cont_points)

# --- Slide 6: LLM Limitations ---
limitation_points = [
    "Not Updated to latest information: Models only know data up to their training cutoff.",
    "Hallucinations: Output can be factually incorrect or nonsensical but look coherent.",
    "Domain-specific accuracy: LLMs often lack specific organizational knowledge (e.g., specific HR policies).",
    "Source Citations: Difficult to verify sources; lack of citations is an ethical concern.",
    "Updates take long training time: Re-training is computationally intensive and resource-heavy."
]
add_text_slide_with_bullets(prs, "LLM's And Its Limitations", limitation_points)

# --- Slide 7: What is RAG? ---
rag_points = [
    "RAG stands for Retrieval-Augmented Generation.",
    "Combines retrieval and generation processes to enhance LLM capabilities.",
    "Retrieves relevant information from a knowledge base/external source.",
    "Uses retrieved info + internal knowledge to generate contextually relevant responses.",
    "Produces higher-quality, context-aware outputs."
]
add_text_slide_with_bullets(prs, "What is RAG?", rag_points)

# --- Slide 8: Why is RAG Important ---
why_rag_points = [
    "LLM Metaphor: An over-enthusiastic employee who is confident but uninformed on current events.",
    "RAG Solution: Redirects LLM to retrieve info from authoritative sources.",
    "Benefits: Greater control over output; users gain insight into generation sources.",
    "Example: Asking about specific recent awards (e.g., Google Cloud Partner 2023) where a standard LLM fails due to date cutoffs."
]
add_text_slide_with_bullets(prs, "Why is RAG Important?", why_rag_points)

# --- Slide 9: RAG Architecture (Section Header) ---
add_slide(prs, 2, "RAG Architecture")

# --- Slide 10: Generalized RAG Approach ---
rag_arch_points = [
    "Framework to mitigate LLM challenges.",
    "[PLACEHOLDER FOR DIAGRAM: Generalized RAG Approach]",
    "Key Components likely in diagram:",
    "- Embedding Model",
    "- Private/Custom Data & External Data Source",
    "- Vector DB / Retrieval Methods",
    "- Retrieved Context",
    "- Query + Prompt -> LLM -> Formatted Response"
]
add_text_slide_with_bullets(prs, "Generalized RAG Approach", rag_arch_points)

# --- Slides 11-20: Images/Diagrams ---
# The text content was empty for these pages, implying they are full-page images or diagrams.
for i in range(11, 21):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    tf = txBox.text_frame
    tf.text = f"[PLACEHOLDER FOR IMAGE/DIAGRAM FROM PAGE {i}]"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

prs.save('RAG_Presentation.pptx')